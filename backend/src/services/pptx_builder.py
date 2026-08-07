"""将 Markdown 报告转为简易 PPTX（依赖 python-pptx）。"""

from __future__ import annotations

import io
import re

from pycore.core import get_logger

logger = get_logger()


def _parse_sections(md: str) -> list[tuple[str, list[str]]]:
    """解析 ## 小节 → (标题, 要点行列表)。"""
    lines = (md or "").splitlines()
    title = "报告"
    sections: list[tuple[str, list[str]]] = []
    current: str | None = None
    bullets: list[str] = []

    def flush() -> None:
        nonlocal current, bullets
        if current is not None:
            sections.append((current, bullets[:12]))
        current = None
        bullets = []

    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if line.startswith("# ") and not line.startswith("## "):
            title = line[2:].strip() or title
            continue
        if line.startswith("## "):
            flush()
            current = line[3:].strip() or "章节"
            continue
        if current is None:
            current = "摘要"
        cleaned = re.sub(r"^[-*•]\s*", "", line)
        cleaned = re.sub(r"^\d+\.\s*", "", cleaned)
        if cleaned:
            bullets.append(cleaned[:180])
    flush()
    if not sections:
        sections = [(title, ["（暂无结构化章节，请查看 Markdown 全文）"])]
    return sections


def build_report_pptx(title: str, markdown: str) -> bytes | None:
    """成功返回 pptx 字节；未安装 python-pptx 时返回 None。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001
        logger.warning("python-pptx unavailable: {}", exc)
        return None

    sections = _parse_sections(markdown)
    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title[:80]
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = "太平洋金科·智能行政咨询助手 · 报告草稿"
        except Exception:  # noqa: BLE001
            pass

    body_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    for sec_title, bullets in sections[:12]:
        s = prs.slides.add_slide(body_layout)
        if s.shapes.title:
            s.shapes.title.text = sec_title[:60]
        # 找正文框
        body = None
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            if s.shapes.title and shape == s.shapes.title:
                continue
            body = shape
            break
        if body is None:
            continue
        tf = body.text_frame
        tf.clear()
        if not bullets:
            p = tf.paragraphs[0]
            p.text = "（无要点）"
            p.font.size = Pt(18)
            continue
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
